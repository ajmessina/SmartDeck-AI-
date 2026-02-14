from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from services.themes import get_theme
import os
import math

class PPTXBuilder:
    """
    Elite Presentation Builder - McKinsey/Apple/TED Caliber
    
    Supports multiple slide layouts and design themes.
    """
    
    # Slide dimensions (widescreen 16:9)
    SLIDE_W = 13.333
    SLIDE_H = 7.5

    def _load_theme(self, theme_id: str):
        """Load color palette and fonts from a theme"""
        t = get_theme(theme_id)
        self.COLOR_NAVY = t["dark_bg"]
        self.COLOR_WHITE = t["light_bg"]
        self.COLOR_OFF_WHITE = t["subtle_bg"]
        self.COLOR_SLATE_700 = t["text_primary"]
        self.COLOR_SLATE_500 = t["text_secondary"]
        self.COLOR_SLATE_300 = t["border"]
        self.COLOR_SLATE_100 = t["card_bg"]
        self.COLOR_TEAL = t["accent"]
        self.COLOR_TEAL_DARK = t["accent_dark"]
        self.COLOR_TEAL_LIGHT = t["accent_light"]
        self.COLOR_CORAL = t["danger"]
        self.COLOR_AMBER = t["warning"]
        self.COLOR_BLUE = t["info"]
        self.COLOR_TEXT_ON_DARK = t["text_on_dark"]
        self.FONT_HEADING = t["font_heading"]
        self.FONT_BODY = t["font_body"]

    def build(self, data: dict, output_filename: str, theme_id: str = "corporate_navy") -> str:
        """Build an elite-level presentation from structured data"""
        self._load_theme(theme_id)
        
        prs = Presentation()
        prs.slide_width = Inches(self.SLIDE_W)
        prs.slide_height = Inches(self.SLIDE_H)
        
        slides_data = data.get("slides", [])
        total_slides = len(slides_data)
        
        for i, slide_data in enumerate(slides_data):
            slide_type = slide_data.get("type", "content_slide")
            
            if slide_type == "title_slide":
                self._create_title_slide(prs, slide_data)
            elif slide_type == "executive_summary":
                self._create_executive_summary(prs, slide_data)
            elif slide_type == "metrics_slide":
                self._create_metrics_slide(prs, slide_data)
            elif slide_type == "two_column":
                self._create_two_column_slide(prs, slide_data)
            elif slide_type == "section_divider":
                self._create_section_divider(prs, slide_data)
            elif slide_type == "challenges_slide":
                self._create_challenges_slide(prs, slide_data)
            elif slide_type == "content_slide":
                self._create_content_slide(prs, slide_data, i, total_slides)
            else:
                self._create_content_slide(prs, slide_data, i, total_slides)
        
        # Closing slide
        self._create_closing_slide(prs)
        
        # Save
        output_dir = "generated_pptx"
        os.makedirs(output_dir, exist_ok=True)
        file_path = os.path.join(output_dir, output_filename)
        prs.save(file_path)
        return file_path

    # =========================================================================
    # SHARED HELPERS
    # =========================================================================

    def _add_bg(self, slide, color=None):
        """Add a solid background rectangle"""
        if color is None:
            color = self.COLOR_WHITE
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(self.SLIDE_W), Inches(self.SLIDE_H)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def _add_bar(self, slide, x, y, w, h, color=None):
        """Add a colored bar/rectangle"""
        if color is None:
            color = self.COLOR_TEAL
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        return bar

    def _add_text(self, slide, x, y, w, h, text, font_size=18,
                  color=None, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=None, anchor=None):
        """Add a text box with styling"""
        if color is None:
            color = self.COLOR_SLATE_700
        if font_name is None:
            font_name = self.FONT_BODY
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        if anchor:
            tf.vertical_anchor = anchor
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return box

    def _add_footer(self, slide, slide_num, total):
        """Add footer with slide number and branding"""
        # Thin line
        self._add_bar(slide, 0.8, 7.0, self.SLIDE_W - 1.6, 0.01,
                      self.COLOR_SLATE_300)
        # Page number
        self._add_text(slide, self.SLIDE_W - 1.5, 7.05, 1, 0.35,
                       f"{slide_num}/{total}", font_size=10,
                       color=self.COLOR_SLATE_500,
                       alignment=PP_ALIGN.RIGHT)
        # Branding
        self._add_text(slide, 0.8, 7.05, 2, 0.35,
                       "SmartDeck AI", font_size=10,
                       color=self.COLOR_SLATE_500)

    # =========================================================================
    # 1. TITLE SLIDE - Cinematic Cover
    # =========================================================================

    def _create_title_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Dark navy background
        self._add_bg(slide, self.COLOR_NAVY)
        
        # Geometric accent - large teal circle (bottom-right, cropped)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(self.SLIDE_W - 4), Inches(self.SLIDE_H - 3),
            Inches(6), Inches(6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.COLOR_TEAL_DARK
        circle.line.fill.background()
        
        # Smaller accent circle
        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(self.SLIDE_W - 6), Inches(self.SLIDE_H - 2),
            Inches(3), Inches(3)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = self.COLOR_TEAL
        circle2.line.fill.background()
        
        # Accent bar top
        self._add_bar(slide, 1.2, 1.8, 1.5, 0.06, self.COLOR_TEAL)
        
        # Title
        title = data.get("title", "Executive Presentation")
        self._add_text(slide, 1.2, 2.2, 8, 2.5, title,
                       font_size=44, color=self.COLOR_WHITE,
                       bold=True, font_name=self.FONT_HEADING)
        
        # Subtitle
        subtitle = data.get("subtitle", "AI-Generated Business Intelligence")
        self._add_text(slide, 1.2, 4.8, 7, 0.8, subtitle,
                       font_size=22, color=self.COLOR_SLATE_300)
        
        # Bottom branding line
        self._add_bar(slide, 1.2, 6.2, 2, 0.04, self.COLOR_TEAL)
        self._add_text(slide, 1.2, 6.4, 3, 0.4,
                       "Powered by SmartDeck AI", font_size=11,
                       color=self.COLOR_SLATE_500)

    # =========================================================================
    # 2. EXECUTIVE SUMMARY - Card Layout with Key Metrics
    # =========================================================================

    def _create_executive_summary(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide)
        
        # Left accent bar
        self._add_bar(slide, 0, 0, 0.12, self.SLIDE_H, self.COLOR_TEAL)
        
        # Title
        self._add_text(slide, 0.8, 0.5, 8, 0.7,
                       data.get("title", "Executive Summary"),
                       font_size=32, bold=True, color=self.COLOR_NAVY,
                       font_name=self.FONT_HEADING)
        
        # Teal divider under title
        self._add_bar(slide, 0.8, 1.25, 1.8, 0.04, self.COLOR_TEAL)
        
        bullets = data.get("bullet_points", [])[:6]
        
        if len(bullets) <= 4:
            # Card layout - each bullet gets its own card
            cols = min(len(bullets), 2)
            rows = math.ceil(len(bullets) / cols)
            card_w = 5.4
            card_h = 2.2
            start_x = 0.8
            start_y = 1.8
            gap = 0.3
            
            for idx, bullet in enumerate(bullets):
                col = idx % cols
                row = idx // cols
                cx = start_x + col * (card_w + gap)
                cy = start_y + row * (card_h + gap)
                
                # Card background
                card = self._add_bar(slide, cx, cy, card_w, card_h,
                                     self.COLOR_SLATE_100)
                card.shadow.inherit = False
                
                # Card accent (left edge of card)
                accent_colors = [self.COLOR_TEAL, self.COLOR_BLUE,
                                 self.COLOR_AMBER, self.COLOR_CORAL]
                accent_c = accent_colors[idx % len(accent_colors)]
                self._add_bar(slide, cx, cy, 0.08, card_h, accent_c)
                
                # Number badge
                self._add_text(slide, cx + 0.3, cy + 0.25, 0.5, 0.5,
                               str(idx + 1), font_size=22,
                               bold=True, color=accent_c)
                
                # Bullet text
                self._add_text(slide, cx + 0.9, cy + 0.25, card_w - 1.3,
                               card_h - 0.5, bullet, font_size=15,
                               color=self.COLOR_SLATE_700)
        else:
            # Standard bullet list for many items
            self._render_bullets(slide, bullets, 0.8, 1.8, 10.5, 5.0)
        
        # Footer
        self._add_footer(slide, 2, "")
        
        # Speaker notes
        if "speaker_notes" in data:
            slide.notes_slide.notes_text_frame.text = data["speaker_notes"]

    # =========================================================================
    # 3. CONTENT SLIDE - Clean Hierarchy
    # =========================================================================

    def _create_content_slide(self, prs, data, idx=0, total=0):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide)
        
        # Left accent bar
        self._add_bar(slide, 0, 0, 0.12, self.SLIDE_H, self.COLOR_TEAL)
        
        # Title
        title = data.get("title", "Key Insight")
        self._add_text(slide, 0.8, 0.5, 10, 0.7, title,
                       font_size=30, bold=True, color=self.COLOR_NAVY,
                       font_name=self.FONT_HEADING)
        
        # Divider
        self._add_bar(slide, 0.8, 1.25, 1.5, 0.04, self.COLOR_TEAL)
        
        # Bullets
        bullets = data.get("bullet_points", [])[:6]
        self._render_bullets(slide, bullets, 0.8, 1.7, 10.5, 4.8)
        
        # Footer
        slide_num = len(prs.slides)
        self._add_footer(slide, slide_num, total)
        
        # Speaker notes
        if "speaker_notes" in data:
            slide.notes_slide.notes_text_frame.text = data["speaker_notes"]

    # =========================================================================
    # 4. TWO-COLUMN SLIDE - Side by Side
    # =========================================================================

    def _create_two_column_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide)
        
        # Left accent bar
        self._add_bar(slide, 0, 0, 0.12, self.SLIDE_H, self.COLOR_TEAL)
        
        # Title
        self._add_text(slide, 0.8, 0.5, 10, 0.7,
                       data.get("title", "Comparison"),
                       font_size=30, bold=True, color=self.COLOR_NAVY,
                       font_name=self.FONT_HEADING)
        
        # Divider
        self._add_bar(slide, 0.8, 1.25, 1.5, 0.04, self.COLOR_TEAL)
        
        # Column headers
        left_title = data.get("left_title", "")
        right_title = data.get("right_title", "")
        
        if left_title:
            self._add_text(slide, 0.8, 1.6, 5.5, 0.5, left_title,
                           font_size=20, bold=True, color=self.COLOR_TEAL_DARK)
        if right_title:
            self._add_text(slide, 7, 1.6, 5.5, 0.5, right_title,
                           font_size=20, bold=True, color=self.COLOR_TEAL_DARK)
        
        # Center divider (vertical)
        self._add_bar(slide, 6.5, 1.6, 0.02, 4.8, self.COLOR_SLATE_300)
        
        # Left column
        left_points = data.get("left_points", [])[:5]
        self._render_bullets(slide, left_points, 0.8, 2.2, 5.3, 4.2,
                             font_size=16)
        
        # Right column
        right_points = data.get("right_points", [])[:5]
        self._render_bullets(slide, right_points, 7, 2.2, 5.3, 4.2,
                             font_size=16)
        
        # Footer
        self._add_footer(slide, len(prs.slides), "")
        
        if "speaker_notes" in data:
            slide.notes_slide.notes_text_frame.text = data["speaker_notes"]

    # =========================================================================
    # 5. METRICS SLIDE - Big Numbers
    # =========================================================================

    def _create_metrics_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide)
        
        # Left accent bar
        self._add_bar(slide, 0, 0, 0.12, self.SLIDE_H, self.COLOR_TEAL)
        
        # Title
        self._add_text(slide, 0.8, 0.5, 10, 0.7,
                       data.get("title", "Key Metrics"),
                       font_size=30, bold=True, color=self.COLOR_NAVY,
                       font_name=self.FONT_HEADING)
        
        self._add_bar(slide, 0.8, 1.25, 1.5, 0.04, self.COLOR_TEAL)
        
        metrics = data.get("metrics", [])[:6]
        
        if not metrics:
            # Fallback: use bullet_points as simple content
            bullets = data.get("bullet_points", [])[:6]
            self._render_bullets(slide, bullets, 0.8, 1.7, 10.5, 4.8)
        else:
            cols = min(len(metrics), 3)
            rows = math.ceil(len(metrics) / cols)
            card_w = (self.SLIDE_W - 2.0) / cols - 0.3
            card_h = 2.3
            
            for idx, metric in enumerate(metrics):
                col = idx % cols
                row = idx // cols
                cx = 0.8 + col * (card_w + 0.3)
                cy = 1.7 + row * (card_h + 0.3)
                
                # Card bg
                self._add_bar(slide, cx, cy, card_w, card_h,
                              self.COLOR_SLATE_100)
                
                # Top accent
                accent_colors = [self.COLOR_TEAL, self.COLOR_BLUE,
                                 self.COLOR_AMBER, self.COLOR_CORAL,
                                 self.COLOR_TEAL_DARK, self.COLOR_BLUE]
                self._add_bar(slide, cx, cy, card_w, 0.06,
                              accent_colors[idx % len(accent_colors)])
                
                # Value (big number)
                value = metric.get("value", "")
                self._add_text(slide, cx + 0.3, cy + 0.3, card_w - 0.6, 1,
                               value, font_size=36, bold=True,
                               color=self.COLOR_NAVY)
                
                # Label
                label = metric.get("label", "")
                self._add_text(slide, cx + 0.3, cy + 1.3, card_w - 0.6, 0.5,
                               label, font_size=13,
                               color=self.COLOR_SLATE_500)
                
                # Change indicator
                change = metric.get("change", "")
                if change:
                    change_color = (self.COLOR_TEAL if "+" in change
                                    or "up" in change.lower()
                                    else self.COLOR_CORAL)
                    self._add_text(slide, cx + 0.3, cy + 1.7,
                                   card_w - 0.6, 0.4, change,
                                   font_size=12, color=change_color,
                                   bold=True)
        
        self._add_footer(slide, len(prs.slides), "")
        
        if "speaker_notes" in data:
            slide.notes_slide.notes_text_frame.text = data["speaker_notes"]

    # =========================================================================
    # 6. SECTION DIVIDER - Bold Transition
    # =========================================================================

    def _create_section_divider(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Navy background
        self._add_bg(slide, self.COLOR_NAVY)
        
        # Large geometric accent
        self._add_bar(slide, 0, 3.2, self.SLIDE_W, 0.08, self.COLOR_TEAL)
        
        # Section title
        title = data.get("title", "Next Section")
        self._add_text(slide, 1.5, 2.0, 10, 1.2, title,
                       font_size=44, bold=True, color=self.COLOR_WHITE,
                       font_name=self.FONT_HEADING,
                       anchor=MSO_ANCHOR.BOTTOM)
        
        # Subtitle
        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_text(slide, 1.5, 3.6, 10, 0.8, subtitle,
                           font_size=20, color=self.COLOR_SLATE_300)

    # =========================================================================
    # 7. CHALLENGES SLIDE - Risk Highlights
    # =========================================================================

    def _create_challenges_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide)
        
        # Left accent bar - coral for warning
        self._add_bar(slide, 0, 0, 0.12, self.SLIDE_H, self.COLOR_CORAL)
        
        # Title with warning color
        self._add_text(slide, 0.8, 0.5, 10, 0.7,
                       data.get("title", "Key Challenges"),
                       font_size=30, bold=True, color=self.COLOR_NAVY,
                       font_name=self.FONT_HEADING)
        
        # Coral divider
        self._add_bar(slide, 0.8, 1.25, 1.5, 0.04, self.COLOR_CORAL)
        
        bullets = data.get("bullet_points", [])[:6]
        
        # Each challenge as a row with icon indicator
        y_pos = 1.7
        for idx, bullet in enumerate(bullets):
            # Warning indicator dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.0), Inches(y_pos + 0.12),
                Inches(0.18), Inches(0.18)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.COLOR_CORAL
            dot.line.fill.background()
            
            # Text
            self._add_text(slide, 1.4, y_pos, 10, 0.6, bullet,
                           font_size=17, color=self.COLOR_SLATE_700)
            
            y_pos += 0.85
        
        self._add_footer(slide, len(prs.slides), "")
        
        if "speaker_notes" in data:
            slide.notes_slide.notes_text_frame.text = data["speaker_notes"]

    # =========================================================================
    # 8. CLOSING SLIDE
    # =========================================================================

    def _create_closing_slide(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Navy background
        self._add_bg(slide, self.COLOR_NAVY)
        
        # Geometric accent
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(self.SLIDE_W - 4), Inches(self.SLIDE_H - 3.5),
            Inches(6), Inches(6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.COLOR_TEAL_DARK
        circle.line.fill.background()
        
        # Teal bar
        self._add_bar(slide, 1.5, 3.0, 2, 0.06, self.COLOR_TEAL)
        
        # Thank You
        self._add_text(slide, 1.5, 3.3, 8, 1.2,
                       "Thank You", font_size=54, bold=True,
                       color=self.COLOR_WHITE,
                       font_name=self.FONT_HEADING)
        
        # CTA
        self._add_text(slide, 1.5, 4.6, 6, 0.6,
                       "Questions & Discussion", font_size=22,
                       color=self.COLOR_SLATE_300)
        
        # Branding
        self._add_bar(slide, 1.5, 5.8, 1.5, 0.04, self.COLOR_TEAL)
        self._add_text(slide, 1.5, 5.95, 3, 0.4,
                       "Powered by SmartDeck AI", font_size=11,
                       color=self.COLOR_SLATE_500)

    # =========================================================================
    # BULLET RENDERER (handles overflow / auto-sizing)
    # =========================================================================

    def _render_bullets(self, slide, bullets, x, y, w, h,
                        font_size=18):
        """Render bullets with proper sizing to avoid overflow"""
        if not bullets:
            return
        
        # Auto-adjust font size based on content volume
        total_chars = sum(len(b) for b in bullets)
        num_bullets = len(bullets)
        
        # Scale down if too much text
        if total_chars > 400 or num_bullets > 5:
            font_size = min(font_size, 16)
        if total_chars > 600:
            font_size = min(font_size, 14)
        
        # Calculate spacing
        available_h = h
        spacing_before = max(6, min(14, int(available_h * 3)))
        spacing_after = spacing_before
        
        box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Add bullet character
            p.text = f"\u2022  {bullet_text}"
            p.level = 0
            p.font.name = self.FONT_BODY
            p.font.size = Pt(font_size)
            p.font.color.rgb = self.COLOR_SLATE_700
            p.alignment = PP_ALIGN.LEFT
            p.font.bold = False
            
            p.space_before = Pt(spacing_before)
            p.space_after = Pt(spacing_after)
            p.line_spacing = 1.3
